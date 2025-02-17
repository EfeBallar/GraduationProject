import os
import PyPDF2
from docx import Document
from pptx import Presentation
from flask import jsonify
from sentence_transformers import SentenceTransformer
import torch
import faiss
import numpy as np
import pickle
from dotenv import load_dotenv

load_dotenv()
DOC_PATH = os.getenv("DOC_PATH")
V_DB_PATH = os.getenv("V_DB_PATH")

def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list:
    """
    Breaks a string into overlapping text chunks.
   
    Returns:
        list: A list of text chunk strings.
    """
    words = text.split()
    chunks = []
    index = 0
    while index < len(words):
        chunk_words = words[index: index + chunk_size]
        chunks.append(" ".join(chunk_words))
        if index + chunk_size >= len(words):
            break
        index += chunk_size - chunk_overlap
    return chunks

def generate_chunks_from_file(file_path: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list:
    """   
    Returns:
        List[dict]: List of dictionaries containing:
            - 'chunk': Text chunk.
            - 'page': Page/Slide/chunk sequence number.
    """
    chunks = []
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages, start=1):
                    page_text = page.extract_text() or ""
                    if page_text:
                        page_chunks = chunk_text(page_text, chunk_size, chunk_overlap)
                        for chunk in page_chunks:
                            chunks.append({"chunk": chunk, "page": page_num})
        elif ext == ".docx":
            doc = Document(file_path)
            full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            if full_text:
                doc_chunks = chunk_text(full_text, chunk_size, chunk_overlap)
                for i, chunk in enumerate(doc_chunks, start=1):
                    # For files without inherent pages, using sequential chunk numbers.
                    chunks.append({"chunk": chunk, "page": i})
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                slide_text = "\n".join(slide_text)
                if slide_text:
                    slide_chunks = chunk_text(slide_text, chunk_size, chunk_overlap)
                    for chunk in slide_chunks:
                        chunks.append({"chunk": chunk, "page": slide_num})
        else:
            # For txt, md, etc. files
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
                if text:
                    text_chunks = chunk_text(text, chunk_size, chunk_overlap)
                    for i, chunk in enumerate(text_chunks, start=1):
                        chunks.append({"chunk": chunk, "page": i})
    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
    
    return chunks

def process_multiple_files(doc_dir: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list:
    """    
    Returns:
        List[dict]: A list of dictionaries. Each dictionary contains:
            - 'file': File name.
            - 'chunk': The text chunk.
            - 'page': Page/Slide/chunk sequence number.
    """
    chunks_data = []
    
    for file_name in os.listdir(doc_dir):
        file_path = os.path.join(doc_dir, file_name)
        chunks = generate_chunks_from_file(file_path, chunk_size, chunk_overlap)
        for chunk_data in chunks:
            # Here the key 'file' stores the filename.
            chunks_data.append({
                "file": file_name,
                "chunk": chunk_data["chunk"],
                "page": chunk_data["page"]
            })

    return chunks_data

def embed_chunks(chunks_data: list, model_name: str = "all-MiniLM-L6-v2", convert_to_tensor: bool = True) -> list:
    """
    Returns:
        List[dict]: The input list with an added key 'embedding' for each dictionary.
    """
    # Load the pre-trained model
    model = SentenceTransformer(model_name)
    
    # Extract all text chunks for encoding
    texts = [item["chunk"] for item in chunks_data]
    
    # Generate embeddings for all chunks at once
    embeddings = model.encode(texts, convert_to_tensor=convert_to_tensor)
    
    # Attach each embedding to its corresponding dictionary
    for i, item in enumerate(chunks_data):
        item["embedding"] = embeddings[i]
    
    return chunks_data

# def save_chunks_to_faiss(chunks_data: list,
#                          index_file: str = "faiss_index.idx",
#                          metadata_file: str = "metadata.pkl") -> None:
#     """
#     Save embeddings from chunks_data into a FAISS vector index 
#     """
#     embeddings = []
#     metadata = []
#     ids = []
    
#     # Assign a unique id for each item (starting at 0)
#     for i, item in enumerate(chunks_data):
#         # Convert embedding to a NumPy array if it's a PyTorch tensor
#         if hasattr(item["embedding"], "cpu"):
#             emb = item["embedding"].cpu().numpy()
#         else:
#             emb = item["embedding"]
#         embeddings.append(emb)
        
#         # Use i as the unique ID
#         item_id = i
#         ids.append(item_id)
#         # Save metadata with pdf, chunk, page, and id.
#         metadata.append({
#             "file": item["file"],
#             "chunk": item["chunk"],
#             "page": item.get("page", "Unknown"),
#             "id": item_id
#         })
    
#     # Stack embeddings into a 2D NumPy array of type float32
#     embeddings_matrix = np.vstack(embeddings).astype("float32")
    
#     # Normalize embeddings for cosine similarity (if using inner product on normalized vectors)
#     faiss.normalize_L2(embeddings_matrix)
    
#     embedding_dim = embeddings_matrix.shape[1]
    
#     # Create an index that supports deletion: wrap IndexFlatIP in an IndexIDMap
#     index_flat = faiss.IndexFlatIP(embedding_dim)
#     index = faiss.IndexIDMap(index_flat)
    
#     # Convert ids list to a NumPy array of type int64
#     ids_np = np.array(ids, dtype=np.int64)
    
#     # Add embeddings with associated ids to the index
#     index.add_with_ids(embeddings_matrix, ids_np)
    
#     # Save the FAISS index and metadata
#     faiss.write_index(index, index_file)
#     with open(metadata_file, "wb") as f:
#         pickle.dump(metadata, f)
    
#     print(f"Saved FAISS index with {index.ntotal} vectors to '{index_file}'.")

def save_chunks_to_faiss(chunks_data: list,
                         index_file: str = "faiss_index.idx",
                         metadata_file: str = "metadata.pkl") -> None:
    """
    Save embeddings from chunks_data into a FAISS vector index, creating necessary files if they don't exist.
    """

    index_file = os.path.join(os.getcwd(), index_file)
    metadata_file = os.path.join(os.getcwd(), metadata_file)

    os.makedirs(os.path.dirname(index_file), exist_ok=True)  # Ensure directory exists

    embeddings = []
    metadata = []
    ids = []

    # Check if FAISS index and metadata file exist
    if os.path.exists(index_file):
        index = faiss.read_index(index_file)
    else:
        index = None  # Will create a new one

    if os.path.exists(metadata_file):
        with open(metadata_file, "rb") as f:
            metadata = pickle.load(f)
    else:
        metadata = []

    # Determine the starting ID
    current_max_id = max((item["id"] for item in metadata), default=-1)

    for i, item in enumerate(chunks_data):
        try:
            # Convert embedding to NumPy array if it's a PyTorch tensor
            emb = item["embedding"].cpu().numpy() if hasattr(item["embedding"], "cpu") else item["embedding"]
            embeddings.append(emb)

            # Assign a new unique ID
            new_id = current_max_id + i + 1
            ids.append(new_id)

            # Save metadata
            metadata.append({
                "file": item["file"],
                "chunk": item["chunk"],
                "page": item.get("page", "Unknown"),
                "id": new_id
            })
        except Exception as e:
            print(f"Error processing chunk {i}: {e}")

    if not embeddings:
        print("No embeddings found. Nothing to save.")
        return

    # Stack embeddings into a 2D NumPy array (float32)
    embeddings_matrix = np.vstack(embeddings).astype("float32")
    
    # Normalize embeddings for cosine similarity
    faiss.normalize_L2(embeddings_matrix)

    embedding_dim = embeddings_matrix.shape[1]
    ids_np = np.array(ids, dtype=np.int64)

    # Create index if it doesn't exist
    if index is None:
        index_flat = faiss.IndexFlatIP(embedding_dim)
        index = faiss.IndexIDMap(index_flat)

    # Add new embeddings
    index.add_with_ids(embeddings_matrix, ids_np)

    # Save updated FAISS index and metadata
    faiss.write_index(index, index_file)
    with open(metadata_file, "wb") as f:
        pickle.dump(metadata, f)

    print(f"Saved FAISS index with {index.ntotal} vectors to '{index_file}'.")

def add_chunks_to_faiss(new_chunks_data: list,
                        index_file: str = "faiss_index.idx",
                        metadata_file: str = "metadata.pkl") -> None:
    """
    Add new chunk embeddings and metadata to an existing FAISS index and metadata file.
    """


    index_file = f"{os.getcwd()}\\faiss\\{index_file}"
    metadata_file = f"{os.getcwd()}\\faiss\\{metadata_file}"


    # Load the existing FAISS index (assumed to be an IndexIDMap)
    try:
        index = faiss.read_index(index_file)
    except Exception as e:
        print(e)

    # Load existing metadata (or start with an empty list)
    try:
        with open(metadata_file, "rb") as f:
            metadata = pickle.load(f)
    except FileNotFoundError:
        print(f"Metadata file '{metadata_file}' not found. Starting with an empty metadata list.")
        metadata = []
    new_embeddings = []
    new_ids = []
    
    # Determine the starting id for new vectors
    if metadata:
        current_max_id = max(item["id"] for item in metadata)

    else:
        current_max_id = -1

    # Process each new chunk
    for i, item in enumerate(new_chunks_data):
        try:
            if hasattr(item["embedding"], "cpu"):
                emb = item["embedding"].cpu().numpy()
            else:
                emb = item["embedding"]
            new_embeddings.append(emb)

            
            # Assign a new unique id
            new_id = current_max_id + i + 1
            new_ids.append(new_id)
            
            metadata.append({
                "pdf": item["pdf"],
                "chunk": item["chunk"],
                "page": item.get("page", "Unknown"),
                "id": new_id
            })
        except Exception as e:
            print(e)

    
    # Stack new embeddings into a 2D NumPy array (float32)
    new_embeddings_matrix = np.vstack(new_embeddings).astype("float32")
    
    # Normalize new embeddings
    faiss.normalize_L2(new_embeddings_matrix)
    
    new_ids_np = np.array(new_ids, dtype=np.int64)
    
    # Add new embeddings (with their ids) to the index
    index.add_with_ids(new_embeddings_matrix, new_ids_np)
    
    # Save the updated FAISS index and metadata
    faiss.write_index(index, index_file)
    with open(metadata_file, "wb") as f:
        pickle.dump(metadata, f)
    
    print(f"Added {len(new_chunks_data)} new vectors. Total vectors in index: {index.ntotal}.")

def delete_chunks_from_file(file_name: str,
                            index_file: str = "faiss_index.idx",
                            metadata_file: str = "metadata.pkl") -> None:
    """
    Delete all vector chunks corresponding to a given file from the FAISS vector database.
    """

    index_file = f"{os.getcwd()}\\faiss\\{index_file}"
    metadata_file = f"{os.getcwd()}\\faiss\\{metadata_file}"

    # Load the FAISS index
    index = faiss.read_index(index_file)
    
    # Load the existing metadata
    try:
        with open(metadata_file, "rb") as f:
            metadata = pickle.load(f)
    except FileNotFoundError:
        print(f"Metadata file '{metadata_file}' not found.")
        return
    
    ids_to_delete = []
    updated_metadata = []
    
    # Identify vectors (by id) that belong to the specified file.
    for item in metadata:
        if item.get("pdf") == file_name:
            if "id" in item:
                ids_to_delete.append(item["id"])
            else:
                print("Error: Metadata item missing 'id'. Cannot delete without vector IDs.")
                return
        else:
            updated_metadata.append(item)
    
    if not ids_to_delete:
        return jsonify({"error": f"No chunks found for file {file_name}"}), 409
    
    # Convert the list of ids to a NumPy array of type int64
    ids_to_delete = np.array(ids_to_delete, dtype=np.int64)
    
    # Remove the vectors corresponding to these ids from the FAISS index.
    index.remove_ids(ids_to_delete)
    
    # Save the updated index and metadata back to disk
    faiss.write_index(index, index_file)
    with open(metadata_file, "wb") as f:
        pickle.dump(updated_metadata, f)
    
    print(f"Deleted {len(ids_to_delete)} chunks from file '{file_name}'.")

if __name__ == "__main__":
    course_codes = ['CS307']
    for course_code in course_codes:
        doc_directory = f"{DOC_PATH}\\{course_code}" 
        # Process PDFs and generate text chunks
        chunks_data = process_multiple_files(doc_directory, chunk_size=1000, chunk_overlap=200)
        print(f"Extracted {len(chunks_data)} chunks for {course_code}.")
    
        # Generate embeddings for each chunk using Sentence Transformers
        chunks_data = embed_chunks(chunks_data, model_name="all-MiniLM-L6-v2", convert_to_tensor=True)

        # Save the chunks and their embeddings to a FAISS vector database (with metadata)
        save_chunks_to_faiss(chunks_data, index_file= f"{V_DB_PATH}\\{course_code}_faiss_index.idx", metadata_file=f"{V_DB_PATH}\\{course_code}_metadata.pkl")
        print(f"Processed the documents for {course_code}.")